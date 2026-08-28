#!/usr/bin/env python3
"""Generate a portable .pptx export of the Lakelink Fiber RA customer roleplay deck.

The HTML deck (demo-artifacts/customer-roleplay-presentation.html) is the
authoritative, fully-interactive artifact (speaker notes pane, overview grid,
fullscreen, dark mode). This script produces a straightforward linear .pptx
covering the same 20 slides and speaker notes for offline/portable use
(e.g. presenting from a machine without a browser, or emailing ahead of a call).
"""
import os
OUTPUT_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "customer-roleplay-presentation.pptx")
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- Brand palette (from demo-artifacts/B2B Broadband Revenue Assurance.html) ----
BG = RGBColor(0xF5, 0xF6, 0xF5)
SURFACE = RGBColor(0xFF, 0xFF, 0xFF)
SURFACE2 = RGBColor(0xFB, 0xFB, 0xFA)
DIVIDER = RGBColor(0xDF, 0xE1, 0xDF)
TEXT = RGBColor(0x1C, 0x23, 0x21)
TEXT_MUTED = RGBColor(0x5C, 0x66, 0x63)
TEXT_FAINT = RGBColor(0x9A, 0xA1, 0x9E)
PRIMARY = RGBColor(0x0A, 0x6E, 0x6E)
PRIMARY_HL = RGBColor(0xCF, 0xE3, 0xE1)
ERROR = RGBColor(0x9C, 0x2B, 0x2B)
GOLD = RGBColor(0xA9, 0x7C, 0x0B)
GOLD_HL = RGBColor(0xEC, 0xDF, 0xC0)
BRAND = RGBColor(0xFF, 0x5F, 0x46)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]


def add_bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, size=14, bold=False,
                 color=TEXT, align=PP_ALIGN.LEFT, italic=False, font="Calibri",
                 line_spacing=1.15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.line_spacing = line_spacing
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = color
            run.font.name = font
    return box


def add_card(slide, left, top, width, height, title, body, title_color=TEXT,
             fill=SURFACE2, title_size=13, body_size=11):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = 0.06
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = DIVIDER
    shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p0 = tf.paragraphs[0]
    p0.text = title
    p0.runs[0].font.bold = True
    p0.runs[0].font.size = Pt(title_size)
    p0.runs[0].font.color.rgb = title_color
    if body:
        for line in body.split("\n"):
            p = tf.add_paragraph()
            p.text = line if line else " "
            p.runs[0].font.size = Pt(body_size)
            p.runs[0].font.color.rgb = TEXT_MUTED
            p.space_before = Pt(3)
    return shape


def add_pill(slide, left, top, text, fill=GOLD_HL, color=GOLD, size=10):
    w = Inches(0.09 * len(text) + 0.35)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, Inches(0.28))
    shape.adjustments[0] = 0.5
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(size)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = color
    return shape


def new_slide(kicker, title, sub=None, notes=None):
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    # card background
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(0.4),
                                   Inches(12.333), Inches(6.7))
    card.adjustments[0] = 0.03
    card.fill.solid()
    card.fill.fore_color.rgb = SURFACE
    card.line.color.rgb = DIVIDER
    card.line.width = Pt(1)
    card.shadow.inherit = False
    y = Inches(0.65)
    if kicker:
        add_textbox(slide, Inches(0.75), y, Inches(11.8), Inches(0.35), kicker.upper(),
                    size=12, bold=True, color=PRIMARY)
        y += Inches(0.4)
    add_textbox(slide, Inches(0.75), y, Inches(11.8), Inches(0.7), title,
                size=26, bold=True, color=TEXT)
    y += Inches(0.65)
    if sub:
        add_textbox(slide, Inches(0.75), y, Inches(11.8), Inches(0.5), sub,
                    size=14, color=TEXT_MUTED)
        y += Inches(0.5)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide, y


def notes_text(timing=None, tell=None, cues=None, translate=None, objection=None):
    parts = []
    if timing:
        parts.append(f"Timing: {timing}")
    if tell:
        parts.append(f"\nTalk track:\n{tell}")
    if cues:
        parts.append("\nDemo cues:\n" + "\n".join(f"- {c}" for c in cues))
    if translate:
        parts.append(f"\nTranslation:\n{translate}")
    if objection:
        parts.append(f"\nIf challenged:\n{objection}")
    return "\n".join(parts)


# ============================================================================
# SLIDE 1 — TITLE
# ============================================================================
slide = prs.slides.add_slide(blank)
add_bg(slide, SURFACE)
add_textbox(slide, Inches(0), Inches(1.0), SLIDE_W, Inches(0.4),
            "CUSTOMER ROLEPLAY · LUMEN TECHNOLOGIES", size=13, bold=True, color=PRIMARY,
            align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0), Inches(1.55), SLIDE_W, Inches(0.9),
            "Lakelink Fiber Revenue Assurance", size=36, bold=True, color=TEXT,
            align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0), Inches(2.35), SLIDE_W, Inches(0.5),
            "From fragmented audits to governed detection and recovery — built on Databricks",
            size=16, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0), Inches(3.15), SLIDE_W, Inches(1.0),
            "$250M–$312.5M", size=54, bold=True, color=ERROR, align=PP_ALIGN.CENTER,
            font="Consolas")
pill = add_pill(slide, Inches(4.9), Inches(4.15), "ILLUSTRATIVE — BENCHMARK ESTIMATE AT LUMEN SCALE, NOT MEASURED LUMEN LEAKAGE",
                 fill=GOLD_HL, color=GOLD, size=11)
pill.width = Inches(3.5)
pill.left = Inches((13.333 - 3.5) / 2)
add_textbox(slide, Inches(1.0), Inches(4.85), Inches(11.3), Inches(0.6),
            "Prospect: Lumen Technologies — no Lumen data used        Demo operator: Lakelink Fiber (fictional)        Format: 20–30 min scored roleplay",
            size=12, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
slide.notes_slide.notes_text_frame.text = notes_text(
    timing="1:15",
    tell="At Lumen's scale, a two-to-two-and-a-half percent leakage rate represents roughly $250 million to $312 million a year. That is not one broken billing rule — it is value escaping between systems that are each locally correct: sales, contracts, provisioning, billing, finance, and collections. We built a working revenue-assurance solution for Lakelink Fiber, a fictional carrier with realistic telecom data and operating patterns. No Lumen customer data is used.",
    cues=["Dana (executive): value framing throughout.", "Priya (technical): architecture, governance, control points behind every number."],
)

# ============================================================================
# SLIDE 2 — LEAKAGE LIVES BETWEEN SYSTEMS
# ============================================================================
slide, y = new_slide("Business problem · Setup 1 of 7", "The systems are working. The handoffs are not.",
    "Leakage persists because the definition of “revenue” changes by function.",
    notes=notes_text(timing="1:30",
        tell="Leakage persists because the definition of revenue changes by function. Sales sees bookings, operations sees provisioned service, billing sees invoices, and finance sees recognized revenue. Each view can be correct while the handoff between them is wrong. Our design moves those reconciliations into governed controls that run against the same enterprise data foundation.",
        translate="Tell-show-tell transition: I have framed the business problem. Now I will show the integrated path from source records to a prioritized recovery action."))
items = ["1. Contract price differs from billed price", "2. Discounts bypass authorization",
         "3. Revenue recognition timing diverges from policy", "4. Collections risk is discovered too late",
         "5. Documents and structured records disagree", "6. Identity mappings break across CRM, billing, service, and finance"]
add_card(slide, Inches(0.75), y, Inches(5.7), Inches(4.6), "Where the handoff breaks", "\n".join(items), body_size=13)
add_card(slide, Inches(6.65), y, Inches(5.9), Inches(2.0), "Today", "Quarterly spreadsheet reconciliation — by the time finance finds the mismatch, the claim may be aged, disputed, or unrecoverable", body_size=13)
add_card(slide, Inches(6.65), y + Inches(2.2), Inches(5.9), Inches(2.0), "Target", "Governed, repeatable controls on every pipeline run, against the same enterprise data foundation", body_size=13)
add_textbox(slide, Inches(6.65), y + Inches(4.35), Inches(5.9), Inches(0.3), "Salesforce · Oracle ERP · Refinitiv · Ironclad · MDM", size=11, color=TEXT_FAINT)

# ============================================================================
# SLIDE 3 — GOVERNED PATH
# ============================================================================
slide, y = new_slide("Architecture · Setup 2 of 7", "One governed path from signal to recovery",
    "We did not create another isolated revenue-assurance database.",
    notes=notes_text(timing="1:45",
        tell="We did not create another isolated revenue-assurance database. We built on the existing TM Forum SID estate and added deterministic snapshots for the systems where reconciliation gaps originate. Seven transparent SQL controls create a unified leakage register in Unity Catalog.",
        translate="Case activity is transactional, so assignments, notes, and lifecycle changes live in Lakebase rather than being forced into an analytical table. That is a deliberate trade-off."))
flow = [
    ("Source systems", "TM Forum data (cdm_tmforum.tmf_*, read-only) + simulated Salesforce CPQ, Oracle ERP, Refinitiv FX, Ironclad CLM, MDM"),
    ("Governed reconciliation", "Unity Catalog · 7 SQL controls · document intelligence (ai_parse_document/ai_extract) · forecast anomaly detection (ai_forecast)"),
    ("Unified leakage register [MEASURED]", "gold_leakage_summary — ~48,108 demo exceptions · ~$601.5M demo amount at risk · evidence on every row"),
    ("Persona-specific action", "AI/BI dashboard (Dana) · Genie space (Marcus) · RA Exceptions Console (Marcus) · Lakebase case workflow"),
]
fy = y
for title, body in flow:
    add_card(slide, Inches(0.9), fy, Inches(11.0), Inches(0.95), title, body, body_size=12, title_size=13)
    fy += Inches(1.08)

# ============================================================================
# SLIDE 4 — DEMO CONTRACT
# ============================================================================
slide, y = new_slide("Demo contract · Setup 3 of 7", "In the next 16 minutes, three outcomes",
    "The verbal contract before switching to the live demo.",
    notes=notes_text(timing="0:45",
        tell="The demo has three tests. First, can an executive quantify the exposure? Second, can an architect or auditor prove each exception? Third, can an analyst move from detection to accountable recovery?"))
add_card(slide, Inches(0.75), y, Inches(3.75), Inches(1.5), "1 · Quantify", "One executive view of total exposure and concentration")
add_card(slide, Inches(4.65), y, Inches(3.75), Inches(1.5), "2 · Prove", "Evidence, lineage, and detection method for each exception")
add_card(slide, Inches(8.55), y, Inches(3.75), Inches(1.5), "3 · Recover", "Prioritized ownership, notes, and controlled lifecycle movement")
add_card(slide, Inches(0.75), y + Inches(1.7), Inches(11.55), Inches(1.6), "Live surfaces",
         "Lakelink Fiber — Revenue Assurance Command Center (AI/BI)\nRA Exceptions Console (AppKit)\nUnity Catalog + Architecture tab links", body_size=13)

# ============================================================================
# SLIDE 5 — DEMO 1: EXECUTIVE COMMAND CENTER
# ============================================================================
slide, y = new_slide("Live demo 1 of 7 · Show 3:00", "Executive Command Center",
    "The first question an executive asks is how much revenue is at risk and whether that number is getting better.",
    notes=notes_text(timing="3:00",
        tell="We have moved from an unbounded audit question to a governed register that finance can decompose by cause and customer. The value is not the dashboard itself — it is deciding where to put the next hour of recovery effort.",
        cues=["Open “Lakelink Fiber — Revenue Assurance Command Center” dashboard.", "Executive Summary tiles → Leakage by Check Type → Customer Risk Scorecard."],
        translate="Faster prioritization, a shared finance/sales/billing/collections definition, and exposure that can be tracked run over run."))
kpis = [("$601.5M", "Total amount at risk (demo)"), ("48,108", "Open exceptions (demo)"),
        ("$127.5M", "High-severity at risk (demo)"), ("~8,200", "Accounts affected (demo)")]
kx = Inches(0.75)
for num, lbl in kpis:
    c = add_card(slide, kx, y, Inches(2.75), Inches(1.3), num, lbl, title_size=20, body_size=11)
    kx += Inches(2.9)
add_card(slide, Inches(0.75), y + Inches(1.5), Inches(11.55), Inches(1.9), "Demo cues",
         "1. Show Executive Summary KPI tiles.\n2. State clearly: ~$601.5M and ~48K figures are demo-data outputs, not Lumen measurements.\n3. Move to Leakage by Check Type — AR collection risk dominates.\n4. Move briefly to Customer Risk Scorecard.", body_size=12)

# ============================================================================
# SLIDE 6 — DEMO 2: OVERVIEW
# ============================================================================
slide, y = new_slide("Live demo 2 of 7 · Show 1:30", "Overview: the same truth for an analyst",
    "The executive and analyst should not debate different numbers.",
    notes=notes_text(timing="1:30",
        tell="One definition now serves two altitudes. The dashboard funds the decision; this workspace turns the decision into daily work.",
        cues=["Connect KPI totals to the executive dashboard.", "Show root-cause concentration bar chart."]))
bars = [("ar_collection_risk", "~$500M"), ("rev_rec_timing_mismatch", "~$85.7M"), ("unauthorized_discount", "~$13.9M")]
add_card(slide, Inches(0.75), y, Inches(11.55), Inches(2.2), "RA Exceptions Console → Overview",
         "\n".join(f"{k}: {v}" for k, v in bars), body_size=13)
add_textbox(slide, Inches(0.75), y + Inches(2.4), Inches(11.55), Inches(0.5),
            "Point out loading, error, retry, and empty-state behavior without dwelling on UI mechanics.", size=12, color=TEXT_MUTED)

# ============================================================================
# SLIDE 7 — DEMO 3: EXCEPTION QUEUE
# ============================================================================
slide, y = new_slide("Live demo 3 of 7 · Show 2:00", "Exception Queue: prioritize the work",
    "Forty-eight thousand exceptions are not a to-do list. Prioritization is the product.",
    notes=notes_text(timing="2:00",
        tell="We have reduced a portfolio-level number to the next best recovery action. Now the architect's question becomes: is this evidence, or merely an assertion?"))
rows = [("ar_collection_risk", "Meridian Logistics Corp", "High", "$284,000", "New"),
        ("contract_price_mismatch", "Harbor Point Manufacturing", "High", "$171,500", "Investigating"),
        ("unauthorized_discount", "Palisade Health Network", "Medium", "$96,200", "Recovering")]
tbl_shape = slide.shapes.add_table(4, 5, Inches(0.75), y, Inches(11.55), Inches(1.6))
tbl = tbl_shape.table
headers = ["Check type", "Account", "Severity", "Impact $", "Status"]
for c, h in enumerate(headers):
    cell = tbl.cell(0, c)
    cell.text = h
    cell.text_frame.paragraphs[0].runs[0].font.bold = True
    cell.text_frame.paragraphs[0].runs[0].font.size = Pt(11)
for r, row in enumerate(rows, start=1):
    for c, val in enumerate(row):
        cell = tbl.cell(r, c)
        cell.text = val
        cell.text_frame.paragraphs[0].runs[0].font.size = Pt(11)
add_card(slide, Inches(0.75), y + Inches(1.8), Inches(11.55), Inches(1.6), "Demo cues",
         "1. Sort or scan by amount at risk.\n2. Filter to a high-value root cause or severity.\n3. Show keyboard-accessible filters and row action.\n4. Open one exception.", body_size=12)

# ============================================================================
# SLIDE 8 — DEMO 4: EVIDENCE DRAWER
# ============================================================================
slide, y = new_slide("Live demo 4 of 7 · Show 2:30", "Exception evidence: prove every dollar",
    "A revenue-assurance signal is only useful if finance, audit, and engineering can reproduce it.",
    notes=notes_text(timing="2:30",
        tell="Every exception carries enough provenance to move from 'the model says so' to 'this source record violated this named control.' That shortens dispute resolution and builds trust with finance and audit.",
        translate="The gold register retains the check type, source reference, detection method, severity, and amount at risk. Unity Catalog provides the governed analytical boundary and lineage path back through the reconciliation layer."))
add_card(slide, Inches(0.75), y, Inches(5.6), Inches(3.0), "Detection evidence",
         "Check type: contract_price_mismatch\nReference: ContractNumber CN-40218\nDetection method: rule_based\nSource table: salesforce_source.contract_line_item\nKnown-leakage flag: true", body_size=12)
add_card(slide, Inches(6.55), y, Inches(6.0), Inches(3.0), "Demo cues",
         "1. Show reference identifier and customer context.\n2. Show friendly source context and detection method.\n3. Show risk tier or scorecard state.\n4. Explain: Architecture tab exposes physical objects; business surfaces use friendly labels.\n5. If available, open a lineage or catalog link.", body_size=12)

# ============================================================================
# SLIDE 9 — DEMO 5: MY CASES
# ============================================================================
slide, y = new_slide("Live demo 5 of 7 · Show 2:30", "My Cases: close the recovery loop",
    "Detection without recovery is a better report, not a business outcome.",
    notes=notes_text(timing="2:30",
        tell="The control now has an owner, a decision trail, and a recovery state. That is the bridge from amount at risk to recovered and prevented revenue.",
        translate="Executive: this gives leadership accountability. Technical: Lakebase provides transactional updates and relational constraints for case state, while the analytical evidence remains governed in Unity Catalog."))
stages = ["New", "Investigating\n(requires assignee)", "Recovering", "Recovered / Written off\n(terminal)"]
sx = Inches(0.75)
for s in stages:
    add_card(slide, sx, y, Inches(2.75), Inches(1.1), s.split("\n")[0], "\n".join(s.split("\n")[1:]), title_size=13, body_size=10)
    sx += Inches(2.9)
add_card(slide, Inches(0.75), y + Inches(1.3), Inches(11.55), Inches(2.0), "Demo cues",
         "1. Assign the exception to yourself or an analyst.\n2. Add a note describing the next action.\n3. Move New → Investigating → Recovering.\n4. Explain illegal lifecycle jumps are rejected.\n5. Open My cases to show owned work.", body_size=12)

# ============================================================================
# SLIDE 10 — DEMO 6: ARCHITECTURE
# ============================================================================
slide, y = new_slide("Live demo 6 of 7 · Show 3:00", "Architecture: transparent and deployable",
    "We designed this to fit the data estate rather than requiring an eighteen-month replacement program.",
    notes=notes_text(timing="3:00",
        tell="The path is explicit, governed, and deployable as one repository. We reuse the customer's existing model, isolate writes, and add controls incrementally.",
        cues=["Open Architecture tab in the console; walk the seven-step tour at executive speed.", "Open one live workspace or catalog link if stable."]))
steps = [
    ("1 · Lakehouse", "Existing TM Forum source foundation (cdm_tmforum.tmf_*, read-only)"),
    ("2 · Simulated sources", "*_source snapshots keyed to golden customers"),
    ("3 · Reconciliation pipelines", "Lakeflow SDP · 7 silver + 4 gold materialized views"),
    ("4 · Unified gold register", "gold_leakage_summary"),
    ("5 · Serving", "AI/BI dashboard (Dana) + Genie One / RA Genie Agent (Marcus)"),
    ("6 · Agentic apps", "RA Exceptions Console (React/TypeScript) + Lakebase case store"),
    ("7 · Governance", "Unity Catalog — Metric Views · Domains · Glossary · Access"),
    ("Deployment", "Databricks Asset Bundles; one repository; incremental grants"),
]
gx, gy = Inches(0.75), y
for i, (t, b) in enumerate(steps):
    add_card(slide, gx, gy, Inches(2.75), Inches(1.55), t, b, title_size=12, body_size=10)
    gx += Inches(2.9)
    if (i + 1) % 4 == 0:
        gx = Inches(0.75)
        gy += Inches(1.7)

# ============================================================================
# SLIDE 11 — DEMO 7: AI
# ============================================================================
slide, y = new_slide("Live demo 7 of 7 · Show 2:00", "AI where it earns its place",
    "AI is useful here when it expands coverage beyond deterministic joins without weakening provenance.",
    notes=notes_text(timing="2:00",
        tell="AI is not replacing the control framework. It is extending the kinds of evidence the framework can evaluate. The output still lands in named, governed, reviewable surfaces."))
add_card(slide, Inches(0.75), y, Inches(3.75), Inches(1.7), "Document Intelligence", "Compare Ironclad contract/invoice PDFs with structured records using ai_parse_document + ai_extract.", body_size=11)
add_card(slide, Inches(4.65), y, Inches(3.75), Inches(1.7), "Forecast vs. Actuals", "Actual revenue, ai_forecast interval, and budget comparison per month.", body_size=11)
add_card(slide, Inches(8.55), y, Inches(3.75), Inches(1.7), "Genie", "One scripted natural-language question against governed data.", body_size=11)
add_card(slide, Inches(0.75), y + Inches(1.9), Inches(11.55), Inches(1.6), "AI-as-force-multiplier disclosure",
         "AI-assisted development accelerated repository analysis, UX review, test creation, and cross-vendor code review.\nDeterministic quality gates and independent review were retained.\nDemo anomalies are intentionally seeded; not a claim of production model accuracy.", body_size=11)

# ============================================================================
# SLIDE 12 — DECISIONS & TRADE-OFFS
# ============================================================================
slide, y = new_slide("Engineering rationale · Close 1 of 4", "Decisions and trade-offs",
    "Intentional choices, not missing architecture.",
    notes=notes_text(timing="1:30",
        tell="These are intentional choices, not missing architecture. We optimized the first release for explainability, repeatability, and time to value. In production, the highest-priority extensions are live ingestion where latency demands it, standardized identity mapping, and Lakebase-to-Unity-Catalog synchronization."))
decisions = [
    ("Reuse TM Forum foundation", "Faster value, realistic relationships", "Depends on existing data readiness"),
    ("Batch-first controls", "Repeatable, auditable, demo-safe", "Not real-time ingestion"),
    ("SQL-first reconciliation", "Transparent to finance and audit", "Less custom flexibility"),
    ("One RA schema", "Simpler deployment and grants", "Less physical silver/gold separation"),
    ("Direct identity joins", "Fewer persisted objects", "Mapping logic can be duplicated"),
    ("Lakebase for cases", "Correct transactional behavior", "Requires sync for governed recovery KPIs"),
    ("Dashboard + app + Genie", "Right experience per persona", "Cross-surface state must stay aligned"),
]
tbl_shape = slide.shapes.add_table(len(decisions) + 1, 3, Inches(0.75), y, Inches(11.55), Inches(4.4))
tbl = tbl_shape.table
tbl.columns[0].width = Inches(3.6)
tbl.columns[1].width = Inches(4.2)
tbl.columns[2].width = Inches(3.75)
for c, h in enumerate(["Decision", "Why", "Trade-off"]):
    cell = tbl.cell(0, c)
    cell.text = h
    cell.text_frame.paragraphs[0].runs[0].font.bold = True
    cell.text_frame.paragraphs[0].runs[0].font.size = Pt(11)
for r, (d, w, t) in enumerate(decisions, start=1):
    for c, val in enumerate([d, w, t]):
        cell = tbl.cell(r, c)
        cell.text = val
        cell.text_frame.paragraphs[0].runs[0].font.size = Pt(10)

# ============================================================================
# SLIDE 13 — TWO-CONTROL SPRINT
# ============================================================================
slide, y = new_slide("Recommended next step · Close 2 of 4", "Start with two controls and prove the economics",
    "Proposed first-value sprint: 2–4 weeks.",
    notes=notes_text(timing="1:30",
        tell="The ask is deliberately small. Give us the two reconciliations you understand least and read-only access to the relevant records. We will align definitions, run the controls, and put a defensible number on the board."))
add_card(slide, Inches(0.75), y, Inches(5.7), Inches(4.4), "Sprint plan",
         "1. Select two high-value reconciliations\n2. Connect read-only source data\n3. Validate identity and business definitions\n4. Quantify exposure and false-positive rate\n5. Put prioritized exceptions in front of owners\n6. Build the scale-out business case", body_size=13)
add_card(slide, Inches(6.65), y, Inches(5.9), Inches(2.0), "Recommended starting controls",
         "Contract price versus billed price\nAccounts-receivable aging and collection risk", body_size=13)
add_card(slide, Inches(6.65), y + Inches(2.2), Inches(5.9), Inches(2.2), "Success criteria",
         "Defensible dollars identified\nEvidence traceable to source\nNamed operational owners\nAgreed roadmap based on measured value", body_size=13)

# ============================================================================
# SLIDE 14 — VALUE STORY & CLOSE
# ============================================================================
slide, y = new_slide("Value story · Close 3 of 4", "Quantify → Prioritize → Prove → Recover → Prevent", None,
    notes=notes_text(timing="1:45",
        tell="Let me leave you with the business progression. Leakage stops being an opinion and becomes a register. Nobody recovers all of the gross leakage — a defensible planning range is $125 million to $220 million of recovered and prevented revenue over a phased 12-to-24-month program, subject to validation on Lumen data."))
steps5 = [("1 · Quantify", "A governed leakage register replaces competing spreadsheets"),
          ("2 · Prioritize", "Risk, amount, cause, and customer concentrate recovery effort"),
          ("3 · Prove", "Source evidence, detection method, and lineage accelerate trust"),
          ("4 · Recover", "Ownership, notes, and lifecycle convert exposure into action"),
          ("5 · Prevent", "Repeatable controls catch recurrence before quarter close")]
sx, sy = Inches(0.75), y
for i, (t, b) in enumerate(steps5):
    add_card(slide, sx, sy, Inches(3.75), Inches(1.15), t, b, body_size=11)
    sx += Inches(3.9)
    if (i + 1) % 3 == 0:
        sx = Inches(0.75)
        sy += Inches(1.3)
add_card(slide, Inches(0.75), y + Inches(2.5), Inches(5.6), Inches(0.9), "$250M–$312M",
         "Illustrative annual gross leakage exposure at Lumen scale", title_color=ERROR, title_size=22, body_size=10)
add_card(slide, Inches(6.55), y + Inches(2.5), Inches(6.0), Inches(0.9), "$125M–$220M",
         "Defensible planning range: recovered + prevented over 12–24 months", title_color=PRIMARY, title_size=22, body_size=10)
add_textbox(slide, Inches(0.75), y + Inches(3.6), Inches(11.55), Inches(0.6),
            "Dana, does that create a sufficiently clear economic test for sponsorship? Priya, are the evidence, governance, and integration boundaries strong enough to earn a read-only technical discovery?",
            size=13, italic=True, color=TEXT_MUTED)

# ============================================================================
# SLIDE 15 — APPENDIX A: EXECUTIVE OBJECTIONS
# ============================================================================
slide, y = new_slide("Appendix A · Objection handling", "Executive objections", None)
exec_obj = [
    ('"How do I know the value is real?"', "Exec: figures are explicitly illustrative; the decision gate is the first two controls on your data. Technical bridge: every measured exception retains a source reference and named control."),
    ('"This sounds expensive."', "Exec: reuses existing data estate, starts with two controls, bounded first investment. Technical bridge: serverless resources, SQL-first controls, one deployable repository."),
    ('"How quickly do we see value?"', "A quantified first-value result in two to four weeks, not a full production transformation."),
    ('"What if recovery is much lower than the headline?"', "We assume it will be lower. Gross exposure and recoverable value are different metrics — aged claims, disputes, churn reduce backward recovery."),
    ('"Why Databricks rather than another point solution?"', "The leakage already spans the customer's governed data estate; this brings controls, AI, analytics, lineage, and the app to shared data."),
]
ox, oy = Inches(0.75), y
for i, (q, a) in enumerate(exec_obj):
    w = Inches(5.7) if i < 4 else Inches(11.55)
    add_card(slide, ox, oy, w, Inches(1.35), q, a, title_size=12, body_size=10)
    if i < 4:
        if i % 2 == 0:
            ox = Inches(6.65)
        else:
            ox = Inches(0.75)
            oy += Inches(1.5)
    else:
        oy += Inches(1.5)

# ============================================================================
# SLIDE 16 — APPENDIX B: TECHNICAL OBJECTIONS
# ============================================================================
slide, y = new_slide("Appendix B · Objection handling", "Technical objections", None)
tech_obj = [
    ('"How do identities resolve across systems?"', "Explicit joins from source keys to golden TM Forum identities; a production roadmap should standardize mappings as complexity grows."),
    ('"How do you prevent false positives?"', "Named/testable controls, non-negative amounts, constrained enumerations, flagged unparseable docs, deterministic seeded coverage."),
    ('"Where is PII protected?"', "Unity Catalog is the governance boundary with least-privilege grants and masking; deployed masking must be verified in target env."),
    ('"Why Lakebase instead of Delta for cases?"', "Cases need frequent row-level updates and low-latency transactions; trade-off is a sync pattern needed for metric views."),
    ('"Is this streaming?"', "Batch-first for deterministic, auditable reconciliation; streaming is a production extension where latency justifies it."),
    ('"Does Genie have unrestricted data access?"', "No — the standalone Genie space queries governed UC data and inherits its access controls."),
]
tx, ty = Inches(0.75), y
for i, (q, a) in enumerate(tech_obj):
    add_card(slide, tx, ty, Inches(5.7), Inches(1.35), q, a, title_size=11, body_size=10)
    if i % 2 == 0:
        tx = Inches(6.65)
    else:
        tx = Inches(0.75)
        ty += Inches(1.5)

# ============================================================================
# SLIDE 17 — APPENDIX C: IMPLEMENTED VS ROADMAP
# ============================================================================
slide, y = new_slide("Appendix C · Scope honesty", "Implemented versus roadmap", None)
implemented = "Existing TM Forum catalog reuse\nSimulated Salesforce, Oracle, Refinitiv, Ironclad, MDM sources\nSeven reconciliation control surfaces\nUnified gold leakage register\nReconciliation scorecard\nDocument-intelligence checks\nForecast anomaly surface\nExecutive AI/BI dashboard\nStandalone Genie space\nRA Exceptions Console\nLakebase cases and notes\nArchitecture links and bundle deployment contract"
roadmap = "Live production connectors and continuous ingestion\nEmbedded Genie side panel\nDeployed Unity Catalog metric views\nDeployed business Domains, governed tags, glossary Pages\nLakebase-to-Delta sync for recovery-rate KPIs\nAutomated CRM, billing, or collections write-back\nMulti-tenancy, disaster recovery, production SSO design\nProven production ML accuracy on organic distributions"
add_card(slide, Inches(0.75), y, Inches(5.7), Inches(4.5), "Implemented — safe to demonstrate", implemented, title_color=PRIMARY, body_size=11)
add_card(slide, Inches(6.65), y, Inches(5.9), Inches(4.5), "Roadmap — design-only", roadmap, title_color=TEXT_FAINT, body_size=11)

# ============================================================================
# SLIDE 18 — APPENDIX D: AI FORCE MULTIPLIER
# ============================================================================
slide, y = new_slide("Appendix D · How this was built", "AI as a force multiplier",
    "How AI accelerated the build without replacing engineering controls.",
    notes=notes_text(tell="We used AI as a teammate for breadth, speed, and independent challenge. We did not use it as the final authority. The implementer and reviewer were separated, acceptance criteria were explicit, and deterministic gates remained part of the process."))
accel = "Repository-wide requirements synthesis\nUX and accessibility audit\nParallel implementation planning\nTest generation and configuration review\nIndependent cross-vendor code review\nCustomer-story and objection rehearsal"
controls = "Human-defined acceptance contract\nDeterministic typecheck, lint, and server-build gates\nIndependent reviewer separated from implementer\nDocumented environment limitations\nExplicit distinction between implemented and aspirational capabilities"
add_card(slide, Inches(0.75), y, Inches(5.7), Inches(3.4), "AI accelerated", accel, body_size=12)
add_card(slide, Inches(6.65), y, Inches(5.9), Inches(3.4), "Controls retained", controls, body_size=12)

# ============================================================================
# SLIDE 19 — TIMING CARD
# ============================================================================
slide, y = new_slide("Appendix E · Timing card", "24 minutes prepared + 6 minutes Q&A", None)
timing_rows = [
    ("Slides 1–4: setup and contract", "4:00"), ("Executive dashboard", "3:00"),
    ("App Overview and Queue", "3:30"), ("Evidence drawer", "2:30"),
    ("Assignment and lifecycle", "2:30"), ("Architecture", "3:00"),
    ("AI surface", "2:00"), ("Trade-offs and value close", "3:30"),
    ("Prepared total", "24:00"), ("Q&A buffer", "6:00"),
]
tbl_shape = slide.shapes.add_table(len(timing_rows) + 1, 2, Inches(0.75), y, Inches(6.5), Inches(4.6))
tbl = tbl_shape.table
tbl.columns[0].width = Inches(4.8)
tbl.columns[1].width = Inches(1.7)
for c, h in enumerate(["Segment", "Target"]):
    cell = tbl.cell(0, c)
    cell.text = h
    cell.text_frame.paragraphs[0].runs[0].font.bold = True
    cell.text_frame.paragraphs[0].runs[0].font.size = Pt(11)
for r, (seg, t) in enumerate(timing_rows, start=1):
    tbl.cell(r, 0).text = seg
    tbl.cell(r, 1).text = t
    for c in range(2):
        tbl.cell(r, c).text_frame.paragraphs[0].runs[0].font.size = Pt(10)
add_card(slide, Inches(7.5), y, Inches(4.85), Inches(4.6), "20-minute compression",
         "Executive dashboard: 2:00\nOverview and Queue: 2:30\nEvidence drawer: 2:00\nCase lifecycle: 2:00\nArchitecture: 2:00\nAI: 1:00\nClose: 2:00\n\nDo not cut quantification, evidence, or recovery.", body_size=11)

# ============================================================================
# SLIDE 20 — REHEARSAL SCORECARD
# ============================================================================
slide, y = new_slide("Appendix E (cont.) · Rehearsal scorecard", "How this gets scored", None)
score = [
    ("Demo setup", "Name Lumen as audience, Lakelink as demo operator.\nState problem + illustrative economics before tech.\nContract for quantify/prove/recover."),
    ("Tell-show-tell", "State one claim before each screen.\nPerform one or two intentional actions only.\nTranslate the screen into a business outcome."),
    ("Value communication", "Return every capability to recovered/prevented revenue, prioritization, trust, or time to action.\nLabel all estimated/synthetic figures.\nEnd with a bounded two-control sprint."),
    ("Reading the room", "Answer the asker first.\nGive the other stakeholder a one-sentence translation.\nOffer depth rather than forcing it."),
]
qx, qy = Inches(0.75), y
for i, (t, b) in enumerate(score):
    add_card(slide, qx, qy, Inches(5.7), Inches(1.6), t, b, body_size=10)
    if i % 2 == 0:
        qx = Inches(6.65)
    else:
        qx = Inches(0.75)
        qy += Inches(1.75)
add_card(slide, Inches(0.75), qy, Inches(11.8), Inches(1.1), "Professionalism",
         "Never debug on stage for more than ten seconds. Volunteer material limitations before being challenged. Avoid claiming undeployed metric views, tags, embedded Genie, or production ML accuracy. Close with two direct stakeholder questions.", body_size=10)

prs.save(OUTPUT_PATH)
print(f"Saved {len(prs.slides._sldIdLst)} slides to {OUTPUT_PATH}")
